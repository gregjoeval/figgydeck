"""Tests for the PowerPoint builder: aspect-ratio fit + combined deck.

Uses synthetic PNGs so no source PDF or external tooling is required.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from PIL import Image

from figgydeck.pptx import _fit_within, build_combined_pptx, build_pptx

# Slide image box (inches), mirrored from pptx.py.
BOX_W, BOX_H = 11.3, 6.0
EMU_PER_INCH = 914400


def _png(path: Path, w: int, h: int) -> None:
    Image.new("RGB", (w, h), (120, 80, 160)).save(path)


def _figure_manifest(images_dir: Path, specs) -> list[dict]:
    """specs: list of (number, filename, width_px, height_px)."""
    images_dir.mkdir(parents=True, exist_ok=True)
    manifest = []
    for number, fn, w, h in specs:
        _png(images_dir / fn, w, h)
        manifest.append({
            "number": number, "type": "figure", "title": None,
            "caption": f"caption {number}", "page": 1, "image_filename": fn,
        })
    return manifest


def _pictures(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


# --- _fit_within unit tests -------------------------------------------------

def test_fit_within_tall_image_is_height_bound(tmp_path):
    p = tmp_path / "tall.png"
    _png(p, 400, 800)  # aspect 0.5, taller than box (1.883)
    left, top, w, h = _fit_within(p, 1.0, 0.75, BOX_W, BOX_H)
    assert h == pytest.approx(BOX_H)          # bound by height
    assert w == pytest.approx(BOX_H * 0.5)    # 3.0"
    # centered horizontally within the box
    assert left == pytest.approx(1.0 + (BOX_W - w) / 2)
    assert top == pytest.approx(0.75)


def test_fit_within_wide_image_is_width_bound(tmp_path):
    p = tmp_path / "wide.png"
    _png(p, 1600, 400)  # aspect 4.0, wider than box
    left, top, w, h = _fit_within(p, 1.0, 0.75, BOX_W, BOX_H)
    assert w == pytest.approx(BOX_W)
    assert h == pytest.approx(BOX_W / 4.0)
    assert top == pytest.approx(0.75 + (BOX_H - h) / 2)  # centered vertically


def test_fit_within_bad_image_falls_back_to_width_only(tmp_path):
    left, top, w, h = _fit_within(tmp_path / "missing.png", 1.0, 0.75, BOX_W, BOX_H)
    assert h is None
    assert (left, top, w) == (1.0, 0.75, BOX_W)


# --- build_pptx: aspect ratio preserved on the slide ------------------------

def test_build_pptx_preserves_aspect_ratio(tmp_path):
    from pptx import Presentation

    images = tmp_path / "images"
    # a tall figure that the old code would have stretched to 11.3x6
    manifest = _figure_manifest(images, [("1.1", "img-000.png", 400, 800)])
    out = tmp_path / "deck.pptx"
    build_pptx(manifest, images, "Book", "Chapter", out, verbose=False)

    pres = Presentation(str(out))
    # slide 0 is the title; slide 1 holds the figure
    pics = _pictures(pres.slides[1])
    assert len(pics) == 1
    pic = pics[0]
    # aspect ratio matches the source (0.5), not the old 11.3/6 box
    assert pic.width / pic.height == pytest.approx(0.5, rel=1e-3)
    # fully inside the box, nothing cropped
    assert pic.width <= round(BOX_W * EMU_PER_INCH) + 1
    assert pic.height <= round(BOX_H * EMU_PER_INCH) + 1


# --- build_combined_pptx ----------------------------------------------------

def test_build_combined_pptx_merges_all_chapters(tmp_path):
    from pptx import Presentation

    ch1_imgs = tmp_path / "ch1" / "images"
    ch2_imgs = tmp_path / "ch2" / "images"
    m1 = _figure_manifest(ch1_imgs, [
        ("1.1", "img-000.png", 800, 400),
        ("1.2", "img-001.png", 400, 800),
    ])
    # note: same basename img-000.png as ch1 — pptx embeds by content, no clash
    m2 = _figure_manifest(ch2_imgs, [("2.1", "img-000.png", 600, 600)])

    out = tmp_path / "combined.pptx"
    build_combined_pptx(
        [(m1, ch1_imgs, "Chapter 1"), (m2, ch2_imgs, "Chapter 2")],
        "Book", out, verbose=False,
    )

    pres = Presentation(str(out))
    # 1 title slide + 3 figure slides
    assert len(pres.slides) == 1 + 3
