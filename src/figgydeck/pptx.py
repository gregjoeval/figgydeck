"""PowerPoint deck builder (optional output format).

This is a thin minimal builder. The richer pptx layout used in early figgydeck
prototypes is preserved in `examples/build_pptx_styled.py` for reference.
"""

from __future__ import annotations

import io
import json
from pathlib import Path

# Image placement box on each slide (inches): left, top, width, height.
# The box is centered horizontally on the 13.333" wide, 7.5" tall 16:9 slide.
_BOX_LEFT = 1.0
_BOX_TOP = 0.75
_BOX_WIDTH = 11.3
_BOX_HEIGHT = 6.0

# Image optimization: re-encode figures as JPEG, downscaled so the longest side
# is at most _MAX_IMG_PX. Keeps decks small enough to import into size-capped
# tools (Google Slides rejects PowerPoint files over 100 MB) without visible
# quality loss at slide size. ~1600px across an 11" box is ~145 DPI.
_MAX_IMG_PX = 1600
_JPEG_QUALITY = 82


def _fit_dims(
    iw: int,
    ih: int,
    box_left: float,
    box_top: float,
    box_w: float,
    box_h: float,
) -> tuple[float, float, float, float | None]:
    """Pure geometry: fit an iw x ih image inside the box, centered, aspect
    preserved. Returns (left, top, w, h) in inches; h=None signals a degenerate
    size (caller should add_picture with width only)."""
    if iw <= 0 or ih <= 0:
        return box_left, box_top, box_w, None
    img_ar = iw / ih
    box_ar = box_w / box_h
    if img_ar >= box_ar:      # wider than box -> width-bound
        w, h = box_w, box_w / img_ar
    else:                     # taller than box -> height-bound
        w, h = box_h * img_ar, box_h
    left = box_left + (box_w - w) / 2
    top = box_top + (box_h - h) / 2
    return left, top, w, h


def _fit_within(
    img_path: Path,
    box_left: float,
    box_top: float,
    box_w: float,
    box_h: float,
) -> tuple[float, float, float, float | None]:
    """Return (left, top, w, h) in inches: image scaled to fit inside the box
    with native aspect ratio, centered. No cropping, no distortion.

    Fallback: if the image can't be opened (PIL error), returns h=None, meaning
    the caller should add_picture with width only (pptx preserves aspect ratio
    when a single dimension is given) at the box's top-left.
    """
    from PIL import Image  # local import: keep module import-light

    try:
        with Image.open(img_path) as im:
            iw, ih = im.size
    except Exception:
        # fall back to width-fit (aspect preserved by pptx when only width set)
        return box_left, box_top, box_w, None

    return _fit_dims(iw, ih, box_left, box_top, box_w, box_h)


def _prepare_image(img_path: Path, optimize: bool):
    """Return (source, (iw, ih)) for add_picture.

    When `optimize` is True, re-encode the image as a downscaled JPEG in memory
    (source is a BytesIO); otherwise source is the original path string. The
    returned (iw, ih) are the ORIGINAL pixel dims (aspect ratio is unchanged by
    downscaling, so they're correct for fitting). On any failure, falls back to
    the original file.
    """
    from PIL import Image

    if not optimize:
        try:
            with Image.open(img_path) as im:
                return str(img_path), im.size
        except Exception:
            return str(img_path), (0, 0)

    try:
        with Image.open(img_path) as im:
            im.load()
            iw, ih = im.size
            # Flatten transparency/palette to white; JPEG has no alpha.
            if im.mode in ("RGBA", "LA", "P"):
                rgba = im.convert("RGBA")
                bg = Image.new("RGB", im.size, (255, 255, 255))
                bg.paste(rgba, mask=rgba.split()[-1])
                im = bg
            elif im.mode == "1":
                im = im.convert("L")
            elif im.mode not in ("RGB", "L"):
                im = im.convert("RGB")

            longest = max(iw, ih)
            if longest > _MAX_IMG_PX:
                scale = _MAX_IMG_PX / longest
                im = im.resize(
                    (max(1, round(iw * scale)), max(1, round(ih * scale))),
                    Image.LANCZOS,
                )

            buf = io.BytesIO()
            im.save(buf, format="JPEG", quality=_JPEG_QUALITY, optimize=True)
            buf.seek(0)
            return buf, (iw, ih)
    except Exception:
        # If anything goes wrong, embed the original untouched.
        try:
            with Image.open(img_path) as im:
                return str(img_path), im.size
        except Exception:
            return str(img_path), (0, 0)


def _new_presentation():
    """Create a blank 16:9 presentation; return (presentation, blank_layout)."""
    from pptx import Presentation
    from pptx.util import Inches

    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    blank_layout = pres.slide_layouts[6]  # blank
    return pres, blank_layout


def _add_title_slide(pres, blank_layout, title: str, subtitle: str) -> None:
    """Add a title slide with a bold book title and a smaller subtitle line."""
    from pptx.util import Inches, Pt

    slide = pres.slides.add_slide(blank_layout)
    txt = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2.5))
    p = txt.text_frame.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(36)
    p.runs[0].font.bold = True
    p2 = txt.text_frame.add_paragraph()
    p2.text = subtitle
    p2.runs[0].font.size = Pt(24)


def _add_figure_slides(
    pres,
    blank_layout,
    manifest: list[dict],
    images_dir: Path,
    book_title: str,
    chapter_title: str,
    *,
    optimize: bool = True,
) -> tuple[int, int]:
    """Append one slide per figure/table entry. Return (added, skipped).

    Each slide holds the image (fit to the slide, aspect ratio preserved,
    nothing cropped) with the study content in the speaker notes. When
    `optimize` is True, images are re-encoded as downscaled JPEGs to keep the
    deck small (see `_prepare_image`).
    """
    from pptx.util import Inches

    added = 0
    skipped = 0
    for entry in manifest:
        if not entry.get("image_filename"):
            skipped += 1
            continue
        img_path = images_dir / entry["image_filename"]
        if not img_path.exists():
            skipped += 1
            continue

        slide = pres.slides.add_slide(blank_layout)

        # Image — fit inside the box, preserving aspect ratio, centered.
        source, (iw, ih) = _prepare_image(img_path, optimize)
        left, top, w, h = _fit_dims(
            iw, ih, _BOX_LEFT, _BOX_TOP, _BOX_WIDTH, _BOX_HEIGHT
        )
        if h is None:
            slide.shapes.add_picture(
                source, Inches(left), Inches(top), width=Inches(w)
            )
        else:
            slide.shapes.add_picture(
                source,
                Inches(left), Inches(top),
                width=Inches(w), height=Inches(h),
            )
        if hasattr(source, "close"):
            source.close()

        # Speaker notes carry the study content (title, caption, source)
        kind = entry["type"].capitalize()
        number_label = f"Table {entry['number']}" if kind == "Table" else f"Figure {entry['number']}"
        notes_lines = [
            f"{number_label}",
        ]
        if entry.get("title"):
            notes_lines.append(entry["title"])
        if entry.get("caption"):
            notes_lines.append("")
            notes_lines.append(entry["caption"])
        notes_lines.extend([
            "",
            f"Book: {book_title}",
            f"Chapter: {chapter_title}",
            f"Page: {entry['page']}",
        ])
        slide.notes_slide.notes_text_frame.text = "\n".join(notes_lines)
        added += 1

    return added, skipped


def build_pptx(
    manifest: list[dict] | str | Path,
    images_dir: str | Path,
    book_title: str,
    chapter_title: str,
    output_path: str | Path,
    *,
    optimize_images: bool = True,
    verbose: bool = True,
) -> Path:
    """Build a slide deck where each figure/table gets one slide.

    Each slide:
        - Image centered, scaled to fit (aspect ratio preserved, no cropping)
        - Caption text in speaker notes (so the slide face stays minimal
          for visual recall)

    optimize_images: re-encode figures as downscaled JPEGs (default) to keep the
        file small; pass False to embed the original PNGs at full resolution.
    """
    try:
        from pptx import Presentation  # noqa: F401 — fail early if dep missing
    except ImportError as e:
        raise ImportError(
            "build_pptx requires python-pptx. Install with: pip install figgydeck[pptx]"
        ) from e

    if isinstance(manifest, (str, Path)):
        manifest = json.loads(Path(manifest).read_text())

    images_dir = Path(images_dir)
    output_path = Path(output_path)
    log = print if verbose else (lambda *a, **k: None)

    pres, blank_layout = _new_presentation()
    _add_title_slide(pres, blank_layout, book_title, chapter_title)
    added, skipped = _add_figure_slides(
        pres, blank_layout, manifest, images_dir, book_title, chapter_title,
        optimize=optimize_images,
    )

    pres.save(str(output_path))
    log(f"\nWrote: {output_path}")
    log(f"  slides added: {added}")
    log(f"  skipped: {skipped}")
    return output_path


def build_combined_pptx(
    chapters: list[tuple[list[dict], Path, str]],
    book_title: str,
    output_path: str | Path,
    *,
    optimize_images: bool = True,
    verbose: bool = True,
) -> Path:
    """Build one deck merging every chapter's figures/tables.

    Args:
        chapters: list of (manifest, images_dir, chapter_title). Each manifest
            may be the list from `extract_chapter()` or a path to a
            `manifest.json`.
        book_title: Book title shown on the single title slide and in each
            slide's speaker notes.
        output_path: Where to write the combined `.pptx`.
        optimize_images: re-encode figures as downscaled JPEGs (default) to keep
            the file small enough for size-capped importers (e.g. Google Slides
            rejects PowerPoint files over 100 MB); pass False for full-res PNGs.

    One title slide (book title + chapter count), then every chapter's figure
    slides back-to-back (no divider slides).
    """
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError as e:
        raise ImportError(
            "build_combined_pptx requires python-pptx. Install with: pip install figgydeck[pptx]"
        ) from e

    output_path = Path(output_path)
    log = print if verbose else (lambda *a, **k: None)

    pres, blank_layout = _new_presentation()
    subtitle = f"{len(chapters)} chapter{'s' if len(chapters) != 1 else ''}"
    _add_title_slide(pres, blank_layout, book_title, subtitle)

    total_added = 0
    total_skipped = 0
    for manifest, images_dir, chapter_title in chapters:
        if isinstance(manifest, (str, Path)):
            manifest = json.loads(Path(manifest).read_text())
        added, skipped = _add_figure_slides(
            pres, blank_layout, manifest, Path(images_dir), book_title, chapter_title,
            optimize=optimize_images,
        )
        total_added += added
        total_skipped += skipped

    pres.save(str(output_path))
    log(f"\nWrote: {output_path}")
    log(f"  chapters: {len(chapters)}")
    log(f"  slides added: {total_added}")
    log(f"  skipped: {total_skipped}")
    return output_path
